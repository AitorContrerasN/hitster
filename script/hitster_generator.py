
###############################
### LIBRERÍAS A UTILIZAR ######
###############################

import qrcode
from PIL import Image
import os
import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from io import BytesIO
import spotipy
from spotipy.oauth2 import SpotifyClientCredentials



################################################
### DATOS A INCLUIR POR PARTE DEL USUARIO ######
################################################

# ¿EN QUÉ MODO QUIERES UTILIZAR EL SCRIPT?
modo_de_uso = 'desde_excel_completo' # 'desde_excel', 'desde_excel_sin_año', 'desde_excel_completo' o 'desde_playlist'

# SEGÚN MODO, ELIGE EL INPUT DE DATOS
playlist_url = 'https://open.spotify.com/playlist/34YN7RrEkxhtfMJuqCBkhk?si=142760cead8b4007'
excel_url = '/Users/aitor/Desktop/hitster/archivos_para_importar/import_en_excel.xlsx'

# Ruta para guardar los archivos finales de PPT y EXCEL
pptx_path = "/Users/aitor/Desktop/hitster/resultados_exportados/tarjetas_a_imprimir.pptx"
xlsx_path = "/Users/aitor/Desktop/hitster/resultados_exportados/listado_canciones.xlsx"

# Datos personales del usuario para el uso de Spotify
client_id = 'b0ce861df37447f5965272f0f620c964'
client_secret = 'be5e4971b3cf4cfe8bb0b585e77684bd'



################################################
### AUTHENTICACIÓN PARA LA API DE SPOTIFY ######
################################################

# La autenticación hay que hacerla para cualquiera de los dos modos (o desde playlist o desde Excel)
auth_manager = SpotifyClientCredentials(client_id=client_id, client_secret=client_secret)
sp = spotipy.Spotify(auth_manager=auth_manager)


###################################
### FUNCIONES PARA EL SCRIPT ######
###################################

# FUNCIÓN A: para obtener datos de canciones a partir de una URL de una lista de reproducción de Spotify
def obtener_info_playlist(url):
    playlist_id = url.split("/")[-1].split("?")[0]
    results = sp.playlist_tracks(playlist_id)
    
    total_canciones = len(results['items'])  # Total de canciones en la playlist
    canciones = []
    posicion = 1  # Iniciar contador para la posición de la canción
    
    for item in results['items']:
        track = item['track']
        titulo = track['name']
        artistas = ', '.join([artist['name'] for artist in track['artists']])
        url_cancion = track['external_urls']['spotify']
        album = track['album']
        año_publicacion = album['release_date'][:4]  # Asume formato YYYY-MM-DD
        
        # Calcular la proporción de la posición respecto al total de canciones
        proporcion = f'{posicion} / {total_canciones}'
        
        canciones.append({
            'URL': url_cancion,
            'Titulo': titulo,
            'Artista': artistas,
            'Año': año_publicacion,
            'Pos': posicion,  # Posición actual de la canción
            'PosRel': proporcion  # Formato 'POSICION / TOTAL CANCIONES' como string
        })
        posicion += 1  # Actualizar posición para la próxima canción
    
    return canciones


# FUNCIÓN B: para obtener datos de canciones a partir de un archivo de Excel con url de canciones de Spotify
def obtener_info_canciones_desde_excel(ruta_archivo):
    # Cargar datos del archivo Excel
    df_urls = pd.read_excel(ruta_archivo)
    
    # Configurar credenciales para la API de Spotify
    client_id = 'tu_client_id'
    client_secret = 'tu_client_secret'
    client_credentials_manager = SpotifyClientCredentials(client_id=client_id, client_secret=client_secret)
    sp = spotipy.Spotify(client_credentials_manager=client_credentials_manager)
    
    # Lista para almacenar los datos de las canciones
    canciones_info = []
    
    # Número total de canciones
    total_canciones = len(df_urls)
    
    # Iterar sobre las URLs en el DataFrame, utilizando enumerate para obtener la posición
    for posicion, url in enumerate(df_urls['URL'], start=1):
        track_id = url.split("/")[-1]
        try:
            # Obtener datos de la canción
            track = sp.track(track_id)
            track_name = track['name']
            track_url = track['external_urls']['spotify']
            release_year = track['album']['release_date'][:4]
            artist_names = ', '.join([artist['name'] for artist in track['artists']])
            proporcion = f'{posicion} / {total_canciones}'
            
            # Añadir la información a la lista
            canciones_info.append({
                'URL': track_url,
                'Titulo': track_name,
                'Artista': artist_names,
                'Año': release_year,
                'Pos': posicion,
                'PosRel': proporcion
            })
        except Exception as e:
            print(f"No se pudo obtener información para la URL: {url}. Error: {str(e)}")
    
    # Crear un DataFrame con la información recolectada
    df_resultado = pd.DataFrame(canciones_info)
    
    return df_resultado



# FUNCIÓN C:
def obtener_info_canciones_desde_excel2(ruta_archivo):
    # Cargar datos del archivo Excel
    df_urls = pd.read_excel(ruta_archivo)
    print(df_urls)
    print('df_urls')
    
    # Configurar credenciales para la API de Spotify
    client_id = 'tu_client_id'
    client_secret = 'tu_client_secret'
    client_credentials_manager = SpotifyClientCredentials(client_id=client_id, client_secret=client_secret)
    sp = spotipy.Spotify(client_credentials_manager=client_credentials_manager)
    
    # Lista para almacenar los datos de las canciones
    canciones_info = []
    print(canciones_info)
    print('canciones_info')
    
    # Número total de canciones
    total_canciones = len(df_urls)
    print(total_canciones)
    print('total_canciones')
    
    # Iterar sobre las filas del DataFrame para procesar cada URL y el año correspondiente
    for posicion, fila in df_urls.iterrows():
        url = fila['URL']
        print(url)
        print('url')
        year = fila['Año']  # Asumiendo que la columna se llama 'Año'
        print(year)
        print('year')
        track_id = url.split("/")[-1]
        print(track_id)
        print('track_id')
        try:
            # Obtener datos de la canción
            track = sp.track(track_id)
            print(track)
            print('track')
            track_name = track['name']
            print(track_name)
            print('track_name')
            track_url = track['external_urls']['spotify']
            artist_names = ', '.join([artist['name'] for artist in track['artists']])
            proporcion = f'{posicion + 1} / {total_canciones}'
            
            # Añadir la información a la lista
            canciones_info.append({
                'URL': track_url,
                'Titulo': track_name,
                'Artista': artist_names,
                'Año': year,  # Utilizar el año del Excel directamente
                'Pos': posicion + 1,
                'PosRel': proporcion
            })
        except Exception as e:
            print(f"No se pudo obtener información para la URL: {url}. Error: {str(e)}")
    
    # Crear un DataFrame con la información recolectada
    df_resultado = pd.DataFrame(canciones_info)
    
    return df_resultado


# FUNCIÓN C:
def obtener_info_canciones_desde_excel_completo(ruta_archivo):
    info_canciones = pd.read_excel(ruta_archivo)
    return info_canciones


###########################################################################################
### APLICACIÓN DE UNO DE LOS DOS MODOS PARA OBTENER LOS DATOS DE LAS CANCIONES EN DF ######
###########################################################################################


if modo_de_uso == 'desde_excel': 
    info_canciones = obtener_info_canciones_desde_excel(excel_url)
    print(info_canciones)
elif modo_de_uso == 'desde_excel_sin_año': 
    info_canciones = obtener_info_canciones_desde_excel2(excel_url)
    print(info_canciones)
elif modo_de_uso == 'desde_playlist':
    info_canciones_lista = obtener_info_playlist(playlist_url)
    info_canciones = pd.DataFrame(info_canciones_lista)
    print(info_canciones)
elif modo_de_uso == 'desde_excel_completo':
    info_canciones = obtener_info_canciones_desde_excel_completo(excel_url)
    print(info_canciones)
else:
    print("Error")



#####################################################################################################
### CREACIÓN DE QRS, Y CONSTRUCCIÓN Y MAQUETACIÓN DEL PPT FINAL CON LOS DATOS DE LAS CANCIONES ######
#####################################################################################################

# Crear una presentación de PowerPoint
prs = Presentation()

# Establecer el tamaño de la diapositiva (19 cm x 19 cm)
prs.slide_width = Cm(19)
prs.slide_height = Cm(19)

# Tamaño del QR y del cuadro de texto
qr_size = Cm(12)  # Tamaño del código QR
text_box_width = Cm(17)  # Ancho del cuadro de texto

# Espacio después de cada párrafo
space_after = Pt(12)

# Calcular posición centrada de la caja de texto
text_box_left = (prs.slide_width - text_box_width) / 2
text_box_top = (prs.slide_height - qr_size) / 2  # Asumiendo que queremos centrar respecto al QR

for index, row in info_canciones.iterrows():
    # Datos de cada fila
    url = row['URL']
    title = row['Titulo']
    artist = row['Artista']
    year = row['Año']

    # Generar el código QR
    qr = qrcode.QRCode(
        version=1,
        error_correction=qrcode.constants.ERROR_CORRECT_L,
        box_size=10,
        border=4,
    )
    qr.add_data(url)
    qr.make(fit=True)
    img = qr.make_image(fill='black', back_color='white')
    bio = BytesIO()
    img.save(bio)
    bio.seek(0)

    # Añadir diapositiva para el QR, centrado
    slide_qr = prs.slides.add_slide(prs.slide_layouts[5])
    qr_left = (prs.slide_width - qr_size) / 2
    qr_top = (prs.slide_height - qr_size) / 2  # Centrado vertical
    slide_qr.shapes.add_picture(bio, qr_left, qr_top, width=qr_size, height=qr_size)

    # Añadir diapositiva para los datos, centrado
    slide_info = prs.slides.add_slide(prs.slide_layouts[5])
    textbox = slide_info.shapes.add_textbox(text_box_left, text_box_top, text_box_width, qr_size)
    tf = textbox.text_frame
    tf.word_wrap = True

    # Configurar el título
    p = tf.add_paragraph()
    p.text = title
    p.space_after = space_after
    p.alignment = PP_ALIGN.CENTER

    # Agregar salto de línea adicional
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = space_after

    # Configurar el artista
    p = tf.add_paragraph()
    p.text = artist
    p.font.bold = True
    p.space_after = space_after
    p.alignment = PP_ALIGN.CENTER

    # Agregar salto de línea adicional
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = space_after

    # Configurar el año con el tamaño deseado
    p = tf.add_paragraph()
    p.text = str(year)
    p.font.size = Pt(80)  # Tamaño de fuente para el año
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER


#########################################
### EXPORTACIÓN DE LOS RESULTADOS ######
#########################################

# Export DataFrame to an Excel file
info_canciones.to_excel(xlsx_path, index=False)

# Guardar la presentación
prs.save(pptx_path)
print(f"PowerPoint guardado en {pptx_path} con éxito güey!")
